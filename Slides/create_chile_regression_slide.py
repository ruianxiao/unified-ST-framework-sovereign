from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_chile_regression_slide():
    # Create presentation
    prs = Presentation()
    
    # Add a slide with a title layout
    slide_layout = prs.slide_layouts[1]  # Using layout 1 which has a title and content
    slide = prs.slides.add_slide(slide_layout)
    
    # Set slide title
    title = slide.shapes.title
    title.text = "Chile Regression Analysis"
    subtitle = slide.placeholders[1]
    subtitle.text = "Model Performance and Variable Significance"
    
    # Add model overview box
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4.5)
    height = Inches(1.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Model Overview:"
    p = tf.add_paragraph()
    p.text = "• Sample Size: 76 observations"
    p = tf.add_paragraph()
    p.text = "• R-squared: 50.7%"
    p = tf.add_paragraph()
    p.text = "• Adjusted R-squared: 47.2%"
    
    # Add variable significance table
    left = Inches(0.5)
    top = Inches(2.8)
    width = Inches(9)
    height = Inches(2.5)
    
    table = slide.shapes.add_table(6, 4, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(3.0)  # Variable
    table.columns[1].width = Inches(2.0)  # Coefficient
    table.columns[2].width = Inches(2.0)  # P-value
    table.columns[3].width = Inches(2.0)  # Significance
    
    # Set header row
    headers = ['Variable', 'Coefficient', 'P-value', 'Significance']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add data rows
    data = [
        ('FTFXIUSAQ_trans (FX rate)', '0.0021', '0.0026', '***'),
        ('FCPWTI.IUSA_trans_lag2 (Oil price)', '0.0059', '0.0009', '***'),
        ('lnPD_TTC_gap (Price-to-Debt gap)', '0.156', '0.0073', '***'),
        ('FSTOCKPQ_trans (Stock prices)', '-0.000027', '0.0654', '*'),
        ('FLBRQ_trans_lead1 (Labor force)', '0.043', '0.1290', 'ns')
    ]
    
    for i, row in enumerate(data, start=1):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = value
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add model diagnostics box
    left = Inches(5.0)
    top = Inches(1.5)
    width = Inches(4.5)
    height = Inches(1.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Model Diagnostics:"
    p = tf.add_paragraph()
    p.text = "• Heteroskedasticity: p-value = 0.706 (no evidence)"
    p = tf.add_paragraph()
    p.text = "• Durbin-Watson: 1.49 (slight positive autocorrelation)"
    p = tf.add_paragraph()
    p.text = "• VIF values: All < 1.5 (no multicollinearity)"
    
    # Add key findings box
    left = Inches(0.5)
    top = Inches(5.4)
    width = Inches(9)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Key Findings:"
    p = tf.add_paragraph()
    p.text = "• FX rates, oil prices, and price-to-debt gap are the most significant drivers"
    p = tf.add_paragraph()
    p.text = "• Model shows moderate explanatory power (50.7% R-squared)"
    p = tf.add_paragraph()
    p.text = "• All variables are stationary, ensuring reliable statistical inference"
    p = tf.add_paragraph()
    p.text = "• Model includes both contemporaneous and lagged effects"
    
    # Add note box
    left = Inches(0.5)
    top = Inches(7.1)
    width = Inches(9)
    height = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Note: *** p < 0.01, ** p < 0.05, * p < 0.1, ns = not significant"
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.size = Pt(8)
    
    # Create output directory if it doesn't exist
    os.makedirs('Slides/output', exist_ok=True)
    
    # Save the presentation
    prs.save('Slides/output/chile_regression_analysis.pptx')
    print("Slide created successfully in Slides/output/chile_regression_analysis.pptx")

if __name__ == "__main__":
    create_chile_regression_slide() 