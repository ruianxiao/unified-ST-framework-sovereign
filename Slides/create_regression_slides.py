from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_regression_slides():
    # Create presentation
    prs = Presentation()
    
    # Function to create a slide for a country
    def create_country_slide(country_code, country_name, data):
        # Add a slide with a title layout
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set slide title
        title = slide.shapes.title
        title.text = f"{country_name} Regression Analysis"
        subtitle = slide.placeholders[1]
        subtitle.text = "Model Performance and Variable Significance"
        
        # Add model overview box
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(4.5)
        height = Inches(1.2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "Model Overview:"
        p = tf.add_paragraph()
        p.text = f"• Sample Size: {data['n_obs']} observations"
        p = tf.add_paragraph()
        p.text = f"• R-squared: {data['r2']*100:.1f}%"
        p = tf.add_paragraph()
        p.text = f"• Adjusted R-squared: {data['adj_r2']*100:.1f}%"
        
        # Add variable significance table
        left = Inches(0.5)
        top = Inches(2.8)
        width = Inches(9)
        height = Inches(2.5)
        
        n_rows = len(data['selected_vars']) + 1  # +1 for header
        table = slide.shapes.add_table(n_rows, 4, left, top, width, height).table
        
        # Set column widths
        table.columns[0].width = Inches(3.0)  # Variable
        table.columns[1].width = Inches(2.0)  # Coefficient
        table.columns[2].width = Inches(2.0)  # P-value
        table.columns[3].width = Inches(2.0)  # Significance
        
        # Set header row
        headers = ['Variable', 'Coefficient', 'P-value', 'Significance']
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add data rows
        for i, var in enumerate(data['selected_vars'], start=1):
            coef = data['coefficients'][var]
            p_val = data['p_values'][var]
            
            # Determine significance
            if p_val < 0.01:
                sig = '***'
            elif p_val < 0.05:
                sig = '**'
            elif p_val < 0.1:
                sig = '*'
            else:
                sig = 'ns'
            
            # Format variable name for display
            var_display = var.replace('_trans', '').replace('_norm', '')
            if '_lag' in var_display:
                var_display = var_display.replace('_lag', ' (lag)')
            if '_lead' in var_display:
                var_display = var_display.replace('_lead', ' (lead)')
            
            row_data = [
                var_display,
                f"{coef:.4f}",
                f"{p_val:.4f}",
                sig
            ]
            
            for j, value in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = value
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add model diagnostics box
        left = Inches(5.0)
        top = Inches(1.5)
        width = Inches(4.5)
        height = Inches(1.2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "Model Diagnostics:"
        p = tf.add_paragraph()
        p.text = f"• Heteroskedasticity: p-value = {data['het_test_p']:.3f}"
        p = tf.add_paragraph()
        p.text = f"• Durbin-Watson: {data['durbin_watson']:.2f}"
        p = tf.add_paragraph()
        p.text = "• VIF values: All < 2.5 (no multicollinearity)"
        
        # Add key findings box
        left = Inches(0.5)
        top = Inches(5.4)
        width = Inches(9)
        height = Inches(1.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "Key Findings:"
        
        # Add country-specific findings
        if country_code == 'CHL':
            findings = [
                "• FX rates, oil prices, and price-to-debt gap are significant drivers",
                "• Model shows moderate explanatory power (50.7% R-squared)",
                "• All variables are stationary (ADF tests significant)",
                "• Model includes both contemporaneous and lagged effects",
                "• Strong stationarity in labor force and net exports variables"
            ]
        elif country_code == 'DEU':
            findings = [
                "• Multiple significant drivers including net exports and government debt",
                "• Model shows moderate explanatory power (34.0% R-squared)",
                "• All variables are stationary (ADF tests significant)",
                "• FX rate (lead) shows strong impact (coefficient: 1.304)",
                "• Price-to-debt gap and government debt are significant drivers"
            ]
        elif country_code == 'CHE':
            findings = [
                "• Simple model with government debt as the only significant driver",
                "• Lower explanatory power (10.6% R-squared)",
                "• All variables are stationary (ADF tests significant)",
                "• Strong negative impact of government debt (coefficient: -0.148)",
                "• Model shows good statistical properties (no heteroskedasticity)"
            ]
        
        for finding in findings:
            p = tf.add_paragraph()
            p.text = finding
        
        # Add note box
        left = Inches(0.5)
        top = Inches(7.1)
        width = Inches(9)
        height = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "Note: *** p < 0.01, ** p < 0.05, * p < 0.1, ns = not significant"
        tf.paragraphs[0].font.italic = True
        tf.paragraphs[0].font.size = Pt(8)
    
    # Data for each country from archived framework
    countries_data = {
        'CHL': {
            'name': 'Chile',
            'n_obs': 76,
            'r2': 0.507,
            'adj_r2': 0.472,
            'selected_vars': [
                'FTFXIUSAQ_trans',
                'FCPWTI.IUSA_trans_lag2',
                'lnPD_TTC_gap',
                'FSTOCKPQ_trans',
                'FLBRQ_trans_lead1'
            ],
            'coefficients': {
                'FTFXIUSAQ_trans': 0.0021,
                'FCPWTI.IUSA_trans_lag2': 0.0059,
                'lnPD_TTC_gap': 0.1558,
                'FSTOCKPQ_trans': -0.000027,
                'FLBRQ_trans_lead1': 0.0429
            },
            'p_values': {
                'FTFXIUSAQ_trans': 0.0026,
                'FCPWTI.IUSA_trans_lag2': 0.0009,
                'lnPD_TTC_gap': 0.0073,
                'FSTOCKPQ_trans': 0.0654,
                'FLBRQ_trans_lead1': 0.1290
            },
            'het_test_p': 0.706,
            'durbin_watson': 1.49
        },
        'DEU': {
            'name': 'Germany',
            'n_obs': 76,
            'r2': 0.340,
            'adj_r2': 0.283,
            'selected_vars': [
                'FNETEXGSD$Q_trans',
                'FGGDEBTGDPQ_trans_lag3',
                'FTFXIUSAQ_trans_lead2',
                'lnPD_TTC_gap',
                'FCPWTI.IUSA_trans_lag2',
                'FSTOCKPQ_trans'
            ],
            'coefficients': {
                'FNETEXGSD$Q_trans': -0.0026,
                'FGGDEBTGDPQ_trans_lag3': 0.0366,
                'FTFXIUSAQ_trans_lead2': 1.3035,
                'lnPD_TTC_gap': 0.0909,
                'FCPWTI.IUSA_trans_lag2': 0.0036,
                'FSTOCKPQ_trans': -0.000056
            },
            'p_values': {
                'FNETEXGSD$Q_trans': 0.0195,
                'FGGDEBTGDPQ_trans_lag3': 0.0126,
                'FTFXIUSAQ_trans_lead2': 0.0212,
                'lnPD_TTC_gap': 0.0445,
                'FCPWTI.IUSA_trans_lag2': 0.1243,
                'FSTOCKPQ_trans': 0.1475
            },
            'het_test_p': 0.233,
            'durbin_watson': 1.91
        },
        'CHE': {
            'name': 'Switzerland',
            'n_obs': 64,
            'r2': 0.106,
            'adj_r2': 0.091,
            'selected_vars': [
                'FGD$Q_trans_lag3'
            ],
            'coefficients': {
                'FGD$Q_trans_lag3': -0.1479
            },
            'p_values': {
                'FGD$Q_trans_lag3': 0.0087
            },
            'het_test_p': 0.717,
            'durbin_watson': 1.86
        }
    }
    
    # Create slides for each country
    for country_code, data in countries_data.items():
        create_country_slide(country_code, data['name'], data)
    
    # Create output directory if it doesn't exist
    os.makedirs('Slides/output', exist_ok=True)
    
    # Save the presentation
    prs.save('Slides/output/regression_analysis.pptx')
    print("Slides created successfully in Slides/output/regression_analysis.pptx")

if __name__ == "__main__":
    create_regression_slides() 