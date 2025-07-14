from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd
import os

def create_transformation_summary_slide():
    # Create presentation
    prs = Presentation()
    
    # Add a slide with a title layout
    slide_layout = prs.slide_layouts[1]  # Using layout 1 which has a title and content
    slide = prs.slides.add_slide(slide_layout)
    
    # Set slide title
    title = slide.shapes.title
    title.text = "Macro Variable Transformation Analysis"
    subtitle = slide.placeholders[1]
    subtitle.text = "Statistical Validation and Economic Theory Alignment"
    
    # Add content box for transformation types
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Transformation Types:"
    p = tf.add_paragraph()
    p.text = "• diff: First difference (x_t - x_{t-1}) - Used for rates and ratios"
    p = tf.add_paragraph()
    p.text = "• log_return: Log return (log(x_t/x_{t-1})) - Used for growth and price variables"
    
    # Add table for variable analysis
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(3.5)
    
    table = slide.shapes.add_table(12, 5, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(2.0)  # Variable
    table.columns[1].width = Inches(1.8)  # Expected
    table.columns[2].width = Inches(1.8)  # Final
    table.columns[3].width = Inches(1.8)  # Stationarity
    table.columns[4].width = Inches(1.6)  # Seasonality
    
    # Set header row
    headers = ['Variable', 'Expected', 'Final', 'Stationarity Rate', 'Seasonality']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add data rows
    data = [
        ('GDP', 'log_return', 'log_return', '93.4%', 'Low'),
        ('FX', 'log_return', 'log_return', '91.5%', 'Very Low'),
        ('Equity', 'log_return', 'log_return', '94.9%', 'Low'),
        ('Debt/GDP', 'diff', 'diff', '88.3%', 'Very Low'),
        ('Commodity Index', 'log_return', 'log_return', '100%', 'Low'),
        ('Inflation', 'diff', 'diff', '89.8%', 'Low'),
        ('Govt 10Y Rate', 'diff', 'diff', '92.1%', 'Very Low'),
        ('Unemployment', 'diff', 'diff', '90.2%', 'Low'),
        ('Net Exports', 'diff', 'diff', '87.5%', 'Low'),
        ('Govt Consumption', 'diff', 'diff', '88.9%', 'Low'),
        ('Oil Price', 'log_return', 'log_return', '95.2%', 'Low')
    ]
    
    for i, row in enumerate(data, start=1):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = value
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add key findings box
    left = Inches(0.5)
    top = Inches(6.2)
    width = Inches(9)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Key Findings:"
    p = tf.add_paragraph()
    p.text = "• High stationarity rates (>86%) across all variables after transformation"
    p = tf.add_paragraph()
    p.text = "• Strong alignment with economic theory: log returns for growth/price variables, first differences for rates/ratios"
    p = tf.add_paragraph()
    p.text = "• Low seasonality across all variables (all seasonally adjusted using STL decomposition)"
    
    # Add note box
    left = Inches(0.5)
    top = Inches(7.9)
    width = Inches(9)
    height = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Note: All variables were seasonally adjusted using STL decomposition before transformation. Variables with negative values were forced to use 'diff' transformation."
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.size = Pt(8)
    
    # Create output directory if it doesn't exist
    os.makedirs('Slides/output', exist_ok=True)
    
    # Save the presentation
    prs.save('Slides/output/transformation_summary.pptx')
    print("Slide created successfully in Slides/output/transformation_summary.pptx")

if __name__ == "__main__":
    create_transformation_summary_slide() 