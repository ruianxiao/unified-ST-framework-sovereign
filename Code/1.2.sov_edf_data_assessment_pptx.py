import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import datetime
import os
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

# Paths
data_dir = Path('Output/1.2.sov_edf_assessment')
template_path = Path('PPT Template.pptx')
out_pptx = Path('Output/1.2.sov_edf_data_assessment.pptx')

# Load data
gs = pd.read_csv(data_dir/'geo_summary.csv')
sumdf = pd.read_csv(data_dir/'summary.csv')
edf_stats = pd.read_csv(data_dir/'edf_stats.csv')
missing = pd.read_csv(data_dir/'missing_by_country.csv')
suit = pd.read_csv(data_dir/'country_suitability.csv')

# Outlier summary (count per country)
outliers = pd.read_csv(data_dir/'edf_outliers.csv')
outlier_counts = outliers.groupby('cinc')['outlier'].sum().reset_index()
outlier_counts.columns = ['cinc', 'n_outliers']

# Pie/bar chart data
good = (suit['reason'] == 'suitable').sum()
bad = (suit['reason'] != 'suitable').sum()

# Helper: get layouts
prs_template = Presentation(str(template_path))
cover_layout = prs_template.slide_layouts[0]
content_layout = prs_template.slide_layouts[1]
back_cover_1 = prs_template.slides[2]
back_cover_2 = prs_template.slides[3]

MOODYS_BLUE = RGBColor(0, 51, 160)
MOODYS_GRAY = RGBColor(242, 242, 242)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
FONT_NAME = 'Calibri'

# Start from the template and set slide size
prs = Presentation(str(template_path))
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Helper: add title slide
def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(cover_layout)
    for shape in slide.shapes:
        if shape.has_text_frame and 'title' in shape.name.lower():
            shape.text = title
        elif shape.has_text_frame and 'subtitle' in shape.name.lower():
            shape.text = subtitle
    return slide

def set_slide_title(slide, title):
    for shape in slide.shapes:
        if shape.has_text_frame and ('title' in shape.name.lower() or shape.text_frame.text == ''):
            shape.text = title
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(36)
                    run.font.bold = True
                    run.font.name = FONT_NAME
                    run.font.color.rgb = MOODYS_BLUE
            return
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    tf.text = title
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.size = Pt(36)
            run.font.bold = True
            run.font.name = FONT_NAME
            run.font.color.rgb = MOODYS_BLUE

def add_table(slide, df, left, top, width, height, font_size=14):
    rows, cols = df.shape
    table = slide.shapes.add_table(rows+1, cols, left, top, width, height).table
    # Header row
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = MOODYS_BLUE
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(font_size)
                run.font.bold = True
                run.font.name = FONT_NAME
                run.font.color.rgb = WHITE
    # Data rows
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i+1, j)
            cell.text = str(df.iloc[i, j])
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else MOODYS_GRAY
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(font_size)
                    run.font.name = FONT_NAME
                    run.font.color.rgb = BLACK
    return table

def add_paginated_table(prs, df, title, rows_per_slide=15):
    for i in range(0, len(df), rows_per_slide):
        slide = prs.slides.add_slide(content_layout)
        set_slide_title(slide, f'{title} ({i+1}-{min(i+rows_per_slide, len(df))})')
        subdf = df.iloc[i:i+rows_per_slide]
        add_table(slide, subdf, Inches(0.5), Inches(1.0), Inches(8), Inches(4))

def add_image(slide, img_path, left, top, width=None, height=None):
    if width and height:
        slide.shapes.add_picture(str(img_path), left, top, width, height)
    else:
        slide.shapes.add_picture(str(img_path), left, top)

def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(content_layout)
    set_slide_title(slide, title)
    left, top, width, height = Inches(1), Inches(1.5), Inches(11), Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(20)
            run.font.name = FONT_NAME
            run.font.color.rgb = MOODYS_BLUE if i == 0 else BLACK
    return slide

def add_suitability_chart(prs, good, bad):
    import matplotlib.pyplot as plt
    fig, ax = plt.subplots(figsize=(4,3))
    ax.bar(['Suitable', 'Not Suitable'], [good, bad], color=['#0072C6','#D95319'])
    ax.set_ylabel('Number of Countries')
    ax.set_title('Country Suitability for Model')
    chart_path = data_dir/'suitability_bar.png'
    fig.tight_layout()
    plt.savefig(chart_path)
    plt.close(fig)
    slide = prs.slides.add_slide(content_layout)
    set_slide_title(slide, 'Country Suitability Summary')
    add_image(slide, chart_path, Inches(1), Inches(1.5), width=Inches(4), height=Inches(3))
    return slide

def add_country_charts(prs, countries, label):
    n_per_slide = 6
    img_w, img_h = 3, 1.5
    for i in range(0, len(countries), n_per_slide):
        slide = prs.slides.add_slide(content_layout)
        set_slide_title(slide, f'EDF Time Series: {label} ({i+1}-{min(i+n_per_slide, len(countries))})')
        for j, country in enumerate(countries[i:i+n_per_slide]):
            img_path = data_dir/f'edf_timeseries_{country}.png'
            if img_path.exists():
                row, col = divmod(j, 3)
                left = Inches(0.5 + col*3)
                top = Inches(1.0 + row*1.7)
                add_image(slide, img_path, left, top, width=Inches(img_w), height=Inches(img_h))
                txBox = slide.shapes.add_textbox(left, top+Inches(img_h), Inches(img_w), Inches(0.3))
                tf = txBox.text_frame
                tf.text = country

def copy_slide(prs, slide):
    slide_id = prs.slides._next_id
    slide_part = slide.part
    rId = prs.part.relate_to(slide_part, RT.SLIDE)
    prs.slides._sldIdLst.insert_slide(prs.slides._sldIdLst.count, slide_id, rId)

add_title_slide(prs, "Sovereign CDS-Implied EDF Data Assessment", f"Generated: {datetime.date.today()}")

# Data overview
slide = prs.slides.add_slide(content_layout)
set_slide_title(slide, 'Geographic Distribution')
add_table(slide, gs, Inches(0.5), Inches(1.0), Inches(8), Inches(2.5))

# 1Y vs 5Y EDF comparison
slide = prs.slides.add_slide(content_layout)
set_slide_title(slide, '1Y vs 5Y EDF Comparison')
add_image(slide, data_dir/'edf1_vs_edf5_scatter.png', Inches(0.5), Inches(1.0), width=Inches(4), height=Inches(4))
add_table(slide, edf_stats, Inches(5), Inches(1.0), Inches(4), Inches(2.5))
with open(data_dir/'edf1_edf5_corr.txt') as f:
    corr = f.read().strip()
slide.shapes.add_textbox(Inches(5), Inches(4), Inches(4), Inches(0.5)).text = corr

# Data quality assessment
slide = prs.slides.add_slide(content_layout)
set_slide_title(slide, 'Data Quality Assessment')
add_table(slide, missing, Inches(0.5), Inches(1.0), Inches(4), Inches(2.5))
add_table(slide, outlier_counts, Inches(5), Inches(1.0), Inches(4), Inches(2.5))

# Country suitability tables
suitable_table = suit[suit['reason'] == 'suitable'][['cinc','n_quarters','frac_missing','n_extremes','extreme_flag','reason']]
not_suitable_table = suit[suit['reason'] != 'suitable'][['cinc','n_quarters','frac_missing','n_extremes','extreme_flag','reason']]
add_paginated_table(prs, suitable_table, 'Suitable Countries')
add_paginated_table(prs, not_suitable_table, 'Not Suitable Countries')
add_suitability_chart(prs, good, bad)

# Extreme value bar chart
import matplotlib.pyplot as plt
fig, ax = plt.subplots(figsize=(8,3))
suit_sorted = suit.sort_values('n_extremes', ascending=False)
ax.bar(suit_sorted['cinc'], suit_sorted['n_extremes'], color='#D95319')
ax.set_ylabel('Number of Extreme Values')
ax.set_title('Extreme Value Count by Country')
plt.xticks(rotation=90, fontsize=6)
fig.tight_layout()
chart_path = data_dir/'extreme_bar.png'
plt.savefig(chart_path)
plt.close(fig)
slide = prs.slides.add_slide(content_layout)
set_slide_title(slide, 'Extreme Value Assessment')
add_image(slide, chart_path, Inches(0.5), Inches(1.0), width=Inches(8), height=Inches(2.5))

# Per-country EDF time series (grouped, 6 per slide)
suitable_countries = suitable_table['cinc'].tolist()
not_suitable_countries = not_suitable_table['cinc'].tolist()
add_country_charts(prs, suitable_countries, 'Suitable')
add_country_charts(prs, not_suitable_countries, 'Not Suitable')

# Recommendations
bullets = [
    f"Countries analyzed: {sumdf['n_countries'].iloc[0]}",
    f"Time span: {sumdf['date_min'].iloc[0]} to {sumdf['date_max'].iloc[0]} ({sumdf['years'].iloc[0]} years)",
    f"Countries suitable for modeling: {good}",
    f"Countries not suitable: {bad}",
    "Use only countries marked as suitable for model development.",
    "Review countries with high missingness or extreme values before use.",
    "Consider data augmentation or alternative sources for not suitable countries."
]
add_bullet_slide(prs, "Recommendations", bullets)

# Add back cover slides (pg3-4 from template)
# copy_slide(prs, back_cover_1)
# copy_slide(prs, back_cover_2)

prs.save(str(out_pptx))
print(f"PPTX report saved to {out_pptx}") 